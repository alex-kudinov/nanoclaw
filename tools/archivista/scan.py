#!/usr/bin/env python3
"""Archivista multi-source scanner.

Reads source configuration from an Obsidian vault's Sources.md file and
writes file catalog entries as Obsidian markdown pages.

Usage:
    scan.py --sources-file PATH --vault-root PATH [--source SOURCE_ID] [--fresh]
"""

import argparse
import hashlib
import json
import os
import re
import sys
import tempfile
from datetime import datetime, timezone
from pathlib import Path

MAX_FILE_SIZE = 50 * 1024 * 1024
TEXT_PREVIEW_CHARS = 3000

EXTRACTABLE = {
    '.pdf', '.docx', '.doc', '.xlsx', '.xls',
    '.pptx', '.ppt', '.txt', '.md', '.csv',
    '.rtf', '.html', '.htm',
}
SKIP_DIRS = {
    '.git', 'node_modules', '__pycache__',
    '$RECYCLE.BIN', '.Trash', 'System Volume Information',
}


# ── Text extraction ───────────────────────────────────────────────────────────

def extract_text(filepath, ext):
    """Extract text preview from a file. Returns (text, error)."""
    try:
        if ext in {'.txt', '.md', '.csv', '.html', '.htm', '.rtf'}:
            with open(filepath, 'r', errors='replace') as f:
                return f.read(TEXT_PREVIEW_CHARS), None

        if ext == '.pdf':
            return _extract_pdf(filepath)

        if ext in {'.docx', '.doc'}:
            return _extract_docx(filepath)

        if ext in {'.xlsx', '.xls'}:
            return _extract_xlsx(filepath)

        if ext in {'.pptx', '.ppt'}:
            return _extract_pptx(filepath)

        return '', None
    except Exception as e:
        return '', str(e)[:80]


def _extract_pdf(filepath):
    try:
        import pdfplumber
    except ImportError:
        return '', 'pdfplumber not installed'
    with pdfplumber.open(filepath) as pdf:
        text = ''
        for page in pdf.pages[:3]:
            text += (page.extract_text() or '') + '\n'
            if len(text) >= TEXT_PREVIEW_CHARS:
                break
    return text[:TEXT_PREVIEW_CHARS], None


def _extract_docx(filepath):
    try:
        import docx
    except ImportError:
        return '', 'python-docx not installed'
    doc = docx.Document(filepath)
    text = '\n'.join(p.text for p in doc.paragraphs[:50])
    return text[:TEXT_PREVIEW_CHARS], None


def _extract_xlsx(filepath):
    try:
        import openpyxl
    except ImportError:
        return '', 'openpyxl not installed'
    wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
    text = ''
    for sheet in wb.sheetnames[:3]:
        ws = wb[sheet]
        text += f'[{sheet}]\n'
        for row in ws.iter_rows(max_row=20, values_only=True):
            cells = [str(c) if c is not None else '' for c in row]
            text += ' | '.join(cells) + '\n'
            if len(text) >= TEXT_PREVIEW_CHARS:
                break
        if len(text) >= TEXT_PREVIEW_CHARS:
            break
    wb.close()
    return text[:TEXT_PREVIEW_CHARS], None


def _extract_pptx(filepath):
    try:
        from pptx import Presentation
    except ImportError:
        return '', 'python-pptx not installed'
    prs = Presentation(filepath)
    text = ''
    for slide in prs.slides[:10]:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text + '\n'
            if len(text) >= TEXT_PREVIEW_CHARS:
                break
        if len(text) >= TEXT_PREVIEW_CHARS:
            break
    return text[:TEXT_PREVIEW_CHARS], None


# ── Sources.md parsing ────────────────────────────────────────────────────────

def parse_sources(sources_file):
    """Parse Sources.md and return list of active source dicts."""
    path = Path(sources_file).expanduser()
    if not path.exists():
        print(f"Error: sources file not found: {path}", file=sys.stderr)
        sys.exit(1)

    text = path.read_text(encoding='utf-8', errors='replace')
    sources = []
    current_id = None
    current_fields = {}

    # Match `- **Key:** value` (colon inside bold) or `- **Key**: value` (colon outside)
    # Capture group allows colons since Key may appear as "Label:" inside **...**
    kv_re = re.compile(
        r'^\s*-\s+\*\*([^*]+?)\*\*:?\s*(.*)',
        re.IGNORECASE,
    )

    for line in text.splitlines():
        h3 = re.match(r'^###\s+(\S+)', line)
        if h3:
            if current_id is not None:
                sources.append({'id': current_id, **current_fields})
            current_id = h3.group(1)
            current_fields = {}
            continue

        if current_id is None:
            continue

        m = kv_re.match(line)
        if m:
            # Strip trailing colon from key (handles "Label:" inside bold)
            key = m.group(1).strip().rstrip(':').lower()
            val = m.group(2).strip()
            # Strip leading colon from value (handles colon outside bold)
            val = val.lstrip(':').strip()
            current_fields[key] = val

    if current_id is not None:
        sources.append({'id': current_id, **current_fields})

    active = [
        s for s in sources
        if s.get('active', '').lower() == 'yes'
    ]
    return active


# ── Filename sanitization ─────────────────────────────────────────────────────

_UNSAFE_CHARS = re.compile(r'[/\\:*?"<>|]')


def safe_filename(original_name, relpath_from_source):
    """Return a safe .md filename for a catalog entry."""
    stem = Path(original_name).stem
    stem = _UNSAFE_CHARS.sub('_', stem)
    stem = stem.lstrip('.')
    stem = stem[:180] if len(stem) > 180 else stem

    safe = f"{stem}.md"
    return safe


def safe_filename_with_collision(original_name, relpath_from_source, dest_dir):
    """Return safe filename, appending --hash8 if a collision exists."""
    stem = Path(original_name).stem
    stem = _UNSAFE_CHARS.sub('_', stem)
    stem = stem.lstrip('.')
    stem = stem[:180] if len(stem) > 180 else stem

    candidate = dest_dir / f"{stem}.md"
    if not candidate.exists():
        return f"{stem}.md"

    # Collision — check if it's the same logical file (same relpath in frontmatter)
    existing_relpath = _read_frontmatter_field(candidate, 'source-path')
    if existing_relpath == relpath_from_source:
        return f"{stem}.md"

    # Different file — append hash
    hash8 = hashlib.md5(relpath_from_source.encode()).hexdigest()[:8]
    return f"{stem}--{hash8}.md"


# ── Frontmatter helpers ───────────────────────────────────────────────────────

def _read_frontmatter_field(md_path, field):
    """Read a single field from YAML frontmatter. Returns value string or None."""
    try:
        content = md_path.read_text(encoding='utf-8', errors='replace')
    except OSError:
        return None

    m = re.match(r'^---\n(.*?)\n---', content, re.DOTALL)
    if not m:
        return None

    block = m.group(1)
    pattern = re.compile(
        r'^' + re.escape(field) + r':\s*(.*)', re.MULTILINE
    )
    fm = pattern.search(block)
    return fm.group(1).strip() if fm else None


def _read_frontmatter_fields(md_path, fields):
    """Read multiple fields from YAML frontmatter. Returns dict."""
    try:
        content = md_path.read_text(encoding='utf-8', errors='replace')
    except OSError:
        return {}

    m = re.match(r'^---\n(.*?)\n---', content, re.DOTALL)
    if not m:
        return {}

    block = m.group(1)
    result = {}
    for field in fields:
        pattern = re.compile(
            r'^' + re.escape(field) + r':\s*(.*)', re.MULTILINE
        )
        fm = pattern.search(block)
        if fm:
            result[field] = fm.group(1).strip()
    return result


# ── Human-readable size ───────────────────────────────────────────────────────

def human_size(size_bytes):
    if size_bytes < 1024:
        return f"{size_bytes}B"
    kb = size_bytes / 1024
    if kb < 1024:
        return f"{kb:.0f}KB"
    mb = kb / 1024
    return f"{mb:.1f}MB"


# ── Catalog entry writer ──────────────────────────────────────────────────────

def build_catalog_entry(
    source_id, domain, relpath, original_name, ext,
    stat, scanned_at, text_preview,
):
    mtime_iso = datetime.fromtimestamp(
        stat.st_mtime, tz=timezone.utc
    ).isoformat()
    size_str = human_size(stat.st_size)
    ext_no_dot = ext.lstrip('.') if ext else ''

    frontmatter = (
        f"---\n"
        f"type: file-catalog\n"
        f"source: {source_id}\n"
        f"source-path: {relpath}\n"
        f"file-type: {ext_no_dot}\n"
        f"size: {size_str}\n"
        f"source-modified: {mtime_iso}\n"
        f"source-size: {stat.st_size}\n"
        f"scanned: {scanned_at}\n"
        f"domain: {domain}\n"
        f"tags: [{domain}]\n"
        f"---\n"
    )

    preview_section = text_preview.strip() if text_preview else ''
    preview_text = preview_section or "No text content extracted."

    body = (
        f"\n# {original_name}\n\n"
        f"Source: {source_id} / {relpath}\n\n"
        f"## Content Preview\n\n"
        f"{preview_text}\n"
    )

    return frontmatter + body


def write_catalog_entry_atomic(dest_path, content):
    """Write content to dest_path atomically using a temp file in the same dir."""
    dest_dir = dest_path.parent
    dest_dir.mkdir(parents=True, exist_ok=True)

    fd, tmp_path = tempfile.mkstemp(dir=dest_dir, suffix='.tmp')
    try:
        with os.fdopen(fd, 'w', encoding='utf-8') as f:
            f.write(content)
        os.replace(tmp_path, dest_path)
    except Exception:
        try:
            os.unlink(tmp_path)
        except OSError:
            pass
        raise


# ── Scan log ──────────────────────────────────────────────────────────────────

def append_scan_log(vault_root, source_id, scanned, new, updated, errors):
    log_dir = vault_root / 'Archivista'
    log_dir.mkdir(parents=True, exist_ok=True)
    log_path = log_dir / 'Scan Log.md'

    ts = datetime.now().strftime('%Y-%m-%d %H:%M')
    line = (
        f"| {ts} | {source_id} | {scanned} scanned "
        f"| {new} new | {updated} updated | {errors} errors |\n"
    )

    with open(log_path, 'a', encoding='utf-8') as f:
        f.write(line)


# ── Source scanner ────────────────────────────────────────────────────────────

def scan_source(source, vault_root, fresh):
    """Scan a single source. Returns stats dict."""
    source_id = source['id']
    domain = source.get('domain', source_id)
    source_path_raw = source.get('path', '')

    if not source_path_raw:
        print(
            f"Error: source '{source_id}' has no Path defined — skipping",
            file=sys.stderr,
        )
        return {'source': source_id, 'scanned': 0, 'new': 0,
                'updated': 0, 'unchanged': 0, 'errors': 1}

    source_root = Path(source_path_raw).expanduser().resolve()
    if not source_root.exists():
        print(
            f"Error: source '{source_id}' path does not exist: {source_root} — skipping",
            file=sys.stderr,
        )
        return {'source': source_id, 'scanned': 0, 'new': 0,
                'updated': 0, 'unchanged': 0, 'errors': 1}

    dest_dir = vault_root / domain / 'Files'
    dest_dir.mkdir(parents=True, exist_ok=True)

    scanned = new = updated = unchanged = errors = 0
    scanned_at = datetime.now(tz=timezone.utc).isoformat()

    for root, dirs, files in os.walk(source_root):
        dirs[:] = [
            d for d in dirs
            if d not in SKIP_DIRS and not d.startswith('.')
        ]

        for name in files:
            if name.startswith('.'):
                continue

            filepath = Path(root) / name
            relpath = str(filepath.relative_to(source_root))

            try:
                stat = filepath.stat()
            except OSError as e:
                print(
                    f"  stat error: {relpath}: {e}",
                    file=sys.stderr,
                )
                errors += 1
                continue

            if stat.st_size > MAX_FILE_SIZE:
                continue

            ext = filepath.suffix.lower()

            # Determine output filename (handle collisions)
            out_filename = safe_filename_with_collision(name, relpath, dest_dir)
            out_path = dest_dir / out_filename

            # Incremental check
            if not fresh and out_path.exists():
                fm = _read_frontmatter_fields(
                    out_path, ['source-modified', 'source-size']
                )
                existing_mtime = fm.get('source-modified', '')
                existing_size = fm.get('source-size', '')
                file_mtime = datetime.fromtimestamp(
                    stat.st_mtime, tz=timezone.utc
                ).isoformat()
                if (existing_mtime == file_mtime
                        and existing_size == str(stat.st_size)):
                    unchanged += 1
                    scanned += 1
                    if scanned % 100 == 0:
                        print(f"{scanned} files scanned...", file=sys.stderr)
                    continue

            # Extract text
            text_preview = ''
            if ext in EXTRACTABLE:
                text_preview, _err = extract_text(str(filepath), ext)

            # Build and write entry
            content = build_catalog_entry(
                source_id=source_id,
                domain=domain,
                relpath=relpath,
                original_name=name,
                ext=ext,
                stat=stat,
                scanned_at=scanned_at,
                text_preview=text_preview,
            )

            try:
                was_existing = out_path.exists()
                write_catalog_entry_atomic(out_path, content)
                if was_existing:
                    updated += 1
                else:
                    new += 1
            except Exception as e:
                print(
                    f"  write error: {relpath}: {e}",
                    file=sys.stderr,
                )
                errors += 1
                continue

            scanned += 1
            if scanned % 100 == 0:
                print(f"{scanned} files scanned...", file=sys.stderr)

    return {
        'source': source_id,
        'scanned': scanned,
        'new': new,
        'updated': updated,
        'unchanged': unchanged,
        'errors': errors,
    }


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description='Archivista multi-source scanner',
    )
    parser.add_argument(
        '--sources-file', required=True,
        help='Path to Sources.md in vault',
    )
    parser.add_argument(
        '--vault-root', required=True,
        help='Path to vault root',
    )
    parser.add_argument(
        '--source', default=None,
        help='Scan only this source ID (default: all active)',
    )
    parser.add_argument(
        '--fresh', action='store_true',
        help='Discard previous catalog and rescan everything',
    )
    args = parser.parse_args()

    vault_root = Path(args.vault_root).expanduser().resolve()
    if not vault_root.is_dir():
        print(f"Error: vault-root not found: {vault_root}", file=sys.stderr)
        sys.exit(1)

    active_sources = parse_sources(args.sources_file)

    if not active_sources:
        print(
            "Error: no active sources found in Sources.md "
            "(check Active: yes entries)",
            file=sys.stderr,
        )
        sys.exit(1)

    if args.source:
        active_sources = [s for s in active_sources if s['id'] == args.source]
        if not active_sources:
            print(
                f"Error: source '{args.source}' not found or not active",
                file=sys.stderr,
            )
            sys.exit(1)

    for source in active_sources:
        print(
            f"Scanning source: {source['id']} ({source.get('path', '?')})",
            file=sys.stderr,
        )
        stats = scan_source(source, vault_root, args.fresh)

        append_scan_log(
            vault_root,
            stats['source'],
            stats['scanned'],
            stats['new'],
            stats['updated'],
            stats['errors'],
        )

        print(json.dumps(stats))


if __name__ == '__main__':
    main()
