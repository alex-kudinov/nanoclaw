#!/Users/xbohdpukc/dev/NanoClaw/tools/onedrive/.venv/bin/python3
"""OneDrive triage — catalog, classify, and organize files.

Usage:
    python3 triage.py scan <onedrive_path>     # Catalog files + extract text
    python3 triage.py classify                  # Batch classify with sonnet
    python3 triage.py report                    # Generate markdown triage report
    python3 triage.py execute <path> --dry-run  # Preview moves
    python3 triage.py execute <path> --confirm  # Execute moves

Data stored in tools/onedrive/data/ (catalog.jsonl, triage.jsonl, report.md).
"""

import argparse
import json
import os
import subprocess
import sys
from datetime import datetime, timezone
from pathlib import Path

SCRIPT_DIR = Path(__file__).resolve().parent
DATA_DIR = SCRIPT_DIR / "data"
CATALOG_PATH = DATA_DIR / "catalog.jsonl"
TRIAGE_PATH = DATA_DIR / "triage.jsonl"
REPORT_PATH = DATA_DIR / "report.md"
PROMPT_PATH = SCRIPT_DIR / "classify-prompt.md"

MAX_FILE_SIZE = 50 * 1024 * 1024
TEXT_PREVIEW_CHARS = 1000
BATCH_SIZE = 25

EXTRACTABLE = {
    '.pdf', '.docx', '.doc', '.xlsx', '.xls',
    '.pptx', '.ppt', '.txt', '.md', '.csv',
    '.rtf', '.html', '.htm',
}
SKIP_DIRS = {
    '.git', 'node_modules', '__pycache__',
    '$RECYCLE.BIN', '.Trash', 'System Volume Information',
    'Drop',
}


# ── Text extraction ──────────────────────────────────────────────────────────

def extract_text(filepath, ext):
    """Extract text preview from a file. Returns (text, error)."""
    try:
        if ext in {'.txt', '.md', '.csv', '.html', '.htm', '.rtf'}:
            with open(filepath, 'r', errors='replace') as f:
                return f.read(TEXT_PREVIEW_CHARS), None

        if ext == '.pdf':
            return _extract_pdf(filepath)

        if ext in {'.docx', '.doc'}:
            return _extract_docx(filepath)

        if ext in {'.xlsx', '.xls'}:
            return _extract_xlsx(filepath)

        if ext in {'.pptx', '.ppt'}:
            return _extract_pptx(filepath)

        return '', None
    except Exception as e:
        return '', str(e)[:80]


def _extract_pdf(filepath):
    try:
        import pdfplumber
    except ImportError:
        return '', 'pdfplumber not installed'
    with pdfplumber.open(filepath) as pdf:
        text = ''
        for page in pdf.pages[:3]:
            text += (page.extract_text() or '') + '\n'
            if len(text) >= TEXT_PREVIEW_CHARS:
                break
    return text[:TEXT_PREVIEW_CHARS], None


def _extract_docx(filepath):
    try:
        import docx
    except ImportError:
        return '', 'python-docx not installed'
    doc = docx.Document(filepath)
    text = '\n'.join(p.text for p in doc.paragraphs[:50])
    return text[:TEXT_PREVIEW_CHARS], None


def _extract_xlsx(filepath):
    try:
        import openpyxl
    except ImportError:
        return '', 'openpyxl not installed'
    wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
    text = ''
    for sheet in wb.sheetnames[:3]:
        ws = wb[sheet]
        text += f'[{sheet}]\n'
        for row in ws.iter_rows(max_row=20, values_only=True):
            cells = [str(c) if c is not None else '' for c in row]
            text += ' | '.join(cells) + '\n'
            if len(text) >= TEXT_PREVIEW_CHARS:
                break
        if len(text) >= TEXT_PREVIEW_CHARS:
            break
    wb.close()
    return text[:TEXT_PREVIEW_CHARS], None


def _extract_pptx(filepath):
    try:
        from pptx import Presentation
    except ImportError:
        return '', 'python-pptx not installed'
    prs = Presentation(filepath)
    text = ''
    for slide in prs.slides[:10]:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text + '\n'
            if len(text) >= TEXT_PREVIEW_CHARS:
                break
        if len(text) >= TEXT_PREVIEW_CHARS:
            break
    return text[:TEXT_PREVIEW_CHARS], None


# ── Commands ─────────────────────────────────────────────────────────────────

def cmd_scan(args):
    """Walk tree, catalog files, extract text previews."""
    onedrive = Path(args.path).expanduser().resolve()
    if not onedrive.is_dir():
        print(f"Error: {onedrive} is not a directory", file=sys.stderr)
        sys.exit(1)

    DATA_DIR.mkdir(parents=True, exist_ok=True)

    scanned = set()
    if CATALOG_PATH.exists() and not args.fresh:
        with open(CATALOG_PATH) as f:
            for line in f:
                try:
                    scanned.add(json.loads(line)['relpath'])
                except (json.JSONDecodeError, KeyError):
                    pass
        if scanned:
            print(f"Resuming: {len(scanned)} files already cataloged")

    mode = 'w' if args.fresh else 'a'
    total = new = skipped = errors = 0
    ext_counts = {}

    with open(CATALOG_PATH, mode) as out:
        for root, dirs, files in os.walk(onedrive):
            dirs[:] = [d for d in dirs
                       if d not in SKIP_DIRS and not d.startswith('.')]

            for name in files:
                if name.startswith('.'):
                    continue

                total += 1
                filepath = Path(root) / name
                relpath = str(filepath.relative_to(onedrive))

                if relpath in scanned:
                    continue

                ext = filepath.suffix.lower()
                ext_counts[ext] = ext_counts.get(ext, 0) + 1

                try:
                    stat = filepath.stat()
                except OSError:
                    errors += 1
                    continue

                if stat.st_size > MAX_FILE_SIZE:
                    skipped += 1
                    continue

                text, extract_err = '', None
                if ext in EXTRACTABLE:
                    text, extract_err = extract_text(str(filepath), ext)

                record = {
                    'relpath': relpath,
                    'ext': ext,
                    'size': stat.st_size,
                    'mtime': datetime.fromtimestamp(
                        stat.st_mtime, tz=timezone.utc
                    ).isoformat(),
                    'text': text,
                    'extract_error': extract_err,
                }
                out.write(json.dumps(record) + '\n')
                new += 1

                if new % 100 == 0:
                    print(f"  {new} files scanned...", file=sys.stderr)

    print(f"\nScan complete:")
    print(f"  Total found: {total}")
    print(f"  Newly cataloged: {new}")
    print(f"  Previously cataloged: {len(scanned)}")
    print(f"  Skipped (>50MB): {skipped}")
    print(f"  Errors: {errors}")
    print(f"\nFile types (top 20):")
    for ext, count in sorted(ext_counts.items(), key=lambda x: -x[1])[:20]:
        print(f"  {ext or '(none)'}: {count}")


def cmd_classify(args):
    """Batch classify cataloged files with sonnet."""
    if not CATALOG_PATH.exists():
        print("Error: run 'scan' first", file=sys.stderr)
        sys.exit(1)

    if not PROMPT_PATH.exists():
        print(f"Error: {PROMPT_PATH} not found", file=sys.stderr)
        sys.exit(1)

    system_prompt = PROMPT_PATH.read_text()

    catalog = []
    with open(CATALOG_PATH) as f:
        for line in f:
            try:
                catalog.append(json.loads(line))
            except json.JSONDecodeError:
                pass

    classified = set()
    if TRIAGE_PATH.exists() and not args.fresh:
        with open(TRIAGE_PATH) as f:
            for line in f:
                try:
                    classified.add(json.loads(line)['relpath'])
                except (json.JSONDecodeError, KeyError):
                    pass
        if classified:
            print(f"Resuming: {len(classified)} already classified")

    pending = [r for r in catalog if r['relpath'] not in classified]
    if not pending:
        print("All files already classified")
        return

    print(f"Classifying {len(pending)} files in batches of {BATCH_SIZE}...")

    mode = 'w' if args.fresh else 'a'
    ok = fail = 0

    with open(TRIAGE_PATH, mode) as out:
        for i in range(0, len(pending), BATCH_SIZE):
            batch = pending[i:i + BATCH_SIZE]
            batch_num = i // BATCH_SIZE + 1

            lines = []
            for j, rec in enumerate(batch):
                size_kb = rec['size'] / 1024
                size_str = (
                    f"{size_kb:.0f}KB"
                    if size_kb < 1024
                    else f"{size_kb / 1024:.1f}MB"
                )
                entry = (
                    f"[{j}] {rec['relpath']} "
                    f"({size_str}, {rec['mtime'][:10]})"
                )
                if rec.get('text'):
                    preview = rec['text'].replace('\n', ' ')[:500]
                    entry += f"\nContent: {preview}"
                lines.append(entry)

            user_msg = "Classify these files:\n\n" + "\n\n".join(lines)

            try:
                result = subprocess.run(
                    [
                        'claude', '--print',
                        '--model', args.model,
                        '--system-prompt', system_prompt,
                        '--max-turns', '1',
                    ],
                    input=user_msg,
                    capture_output=True,
                    text=True,
                    timeout=180,
                )

                if result.returncode != 0:
                    print(
                        f"  Batch {batch_num}: claude error: "
                        f"{result.stderr[:100]}",
                        file=sys.stderr,
                    )
                    fail += len(batch)
                    continue

                response = result.stdout.strip()
                start = response.find('[')
                end = response.rfind(']') + 1
                if start < 0 or end <= start:
                    print(
                        f"  Batch {batch_num}: no JSON array in response",
                        file=sys.stderr,
                    )
                    fail += len(batch)
                    continue

                classifications = json.loads(response[start:end])
                for cls in classifications:
                    idx = cls.get('index', -1)
                    if 0 <= idx < len(batch):
                        cls['relpath'] = batch[idx]['relpath']
                        cls['size'] = batch[idx]['size']
                        out.write(json.dumps(cls) + '\n')
                        ok += 1

                print(f"  Batch {batch_num}: {len(classifications)} classified")

            except subprocess.TimeoutExpired:
                print(f"  Batch {batch_num}: timeout", file=sys.stderr)
                fail += len(batch)
            except json.JSONDecodeError as e:
                print(
                    f"  Batch {batch_num}: JSON parse error: {e}",
                    file=sys.stderr,
                )
                fail += len(batch)
            except Exception as e:
                print(
                    f"  Batch {batch_num}: error: {e}",
                    file=sys.stderr,
                )
                fail += len(batch)

    print(f"\nClassification: {ok} ok, {fail} failed")


def cmd_report(args):
    """Generate markdown triage report."""
    if not TRIAGE_PATH.exists():
        print("Error: run 'classify' first", file=sys.stderr)
        sys.exit(1)

    records = []
    with open(TRIAGE_PATH) as f:
        for line in f:
            try:
                records.append(json.loads(line))
            except json.JSONDecodeError:
                pass

    by_category = {}
    junk = []
    for r in records:
        if r.get('is_junk'):
            junk.append(r)
        else:
            cat = r.get('category', 'uncategorized')
            by_category.setdefault(cat, []).append(r)

    by_project = {}
    for r in records:
        if not r.get('is_junk'):
            proj = r.get('project') or 'unassigned'
            by_project.setdefault(proj, []).append(r)

    total_size = sum(r.get('size', 0) for r in records)
    junk_size = sum(r.get('size', 0) for r in junk)

    with open(REPORT_PATH, 'w') as f:
        f.write("# OneDrive Triage Report\n\n")
        f.write(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}\n\n")

        f.write("## Summary\n\n")
        f.write(f"- **Total files:** {len(records)}\n")
        f.write(f"- **Total size:** {total_size / 1024 / 1024:.1f} MB\n")
        f.write(
            f"- **Junk:** {len(junk)} files "
            f"({junk_size / 1024 / 1024:.1f} MB)\n"
        )
        f.write(f"- **Keep:** {len(records) - len(junk)}\n")
        f.write(f"- **Categories:** {len(by_category)}\n")
        f.write(f"- **Projects:** {len(by_project)}\n\n")

        f.write("## By Category\n\n")
        for cat in sorted(by_category.keys()):
            files = by_category[cat]
            cat_mb = sum(r.get('size', 0) for r in files) / 1024 / 1024
            f.write(f"### {cat} ({len(files)} files, {cat_mb:.1f} MB)\n\n")
            for r in sorted(files, key=lambda x: x.get('relpath', '')):
                f.write(f"- `{r['relpath']}`\n")
                if r.get('summary'):
                    f.write(f"  {r['summary']}\n")
                if r.get('suggested_path'):
                    f.write(f"  -> `{r['suggested_path']}`\n")
            f.write("\n")

        if junk:
            f.write(f"## Junk ({len(junk)} files)\n\n")
            for r in sorted(junk, key=lambda x: x.get('relpath', '')):
                f.write(
                    f"- `{r['relpath']}` — "
                    f"{r.get('summary', 'junk')}\n"
                )
            f.write("\n")

        f.write("## By Project\n\n")
        for proj in sorted(by_project.keys()):
            files = by_project[proj]
            f.write(f"### {proj} ({len(files)} files)\n\n")
            for r in sorted(files, key=lambda x: x.get('relpath', '')):
                f.write(
                    f"- `{r['relpath']}` [{r.get('category', '?')}]\n"
                )
            f.write("\n")

    print(f"Report: {REPORT_PATH}")
    print(f"  {len(records)} files: {len(records) - len(junk)} keep, "
          f"{len(junk)} junk")


def cmd_execute(args):
    """Move files per approved triage plan."""
    if not TRIAGE_PATH.exists():
        print("Error: run 'classify' and 'report' first", file=sys.stderr)
        sys.exit(1)

    onedrive = Path(args.path).expanduser().resolve()
    if not onedrive.is_dir():
        print(f"Error: {onedrive} is not a directory", file=sys.stderr)
        sys.exit(1)

    records = []
    with open(TRIAGE_PATH) as f:
        for line in f:
            try:
                records.append(json.loads(line))
            except json.JSONDecodeError:
                pass

    junk = [r for r in records if r.get('is_junk')]
    moves = [
        r for r in records
        if not r.get('is_junk') and r.get('suggested_path')
    ]

    if args.dry_run:
        print(f"DRY RUN:")
        print(f"  {len(junk)} junk files -> _junk/")
        print(f"  {len(moves)} files to reorganize")
        for r in moves[:30]:
            print(f"    {r['relpath']}")
            print(f"      -> {r['suggested_path']}")
        if len(moves) > 30:
            print(f"    ... and {len(moves) - 30} more")
        return

    if not args.confirm:
        print(
            "Pass --confirm to execute (destructive). "
            "Use --dry-run to preview.",
            file=sys.stderr,
        )
        sys.exit(1)

    moved = errors = 0

    junk_dir = onedrive / '_junk'
    junk_dir.mkdir(exist_ok=True)
    for r in junk:
        src = onedrive / r['relpath']
        dst = junk_dir / Path(r['relpath']).name
        if src.exists():
            try:
                # Handle name collisions in junk dir
                if dst.exists():
                    stem = dst.stem
                    suffix = dst.suffix
                    n = 1
                    while dst.exists():
                        dst = junk_dir / f"{stem}_{n}{suffix}"
                        n += 1
                dst.parent.mkdir(parents=True, exist_ok=True)
                src.rename(dst)
                moved += 1
            except Exception as e:
                print(f"  Error: {r['relpath']}: {e}", file=sys.stderr)
                errors += 1

    for r in moves:
        src = onedrive / r['relpath']
        dst = onedrive / r['suggested_path']
        if src.exists() and src.resolve() != dst.resolve():
            try:
                dst.parent.mkdir(parents=True, exist_ok=True)
                src.rename(dst)
                moved += 1
            except Exception as e:
                print(f"  Error: {r['relpath']}: {e}", file=sys.stderr)
                errors += 1

    print(f"Done: {moved} moved, {errors} errors")


# ── CLI ──────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description='OneDrive triage tool')
    sub = parser.add_subparsers(dest='command')

    p_scan = sub.add_parser('scan', help='Catalog files + extract text')
    p_scan.add_argument('path', help='OneDrive root path')
    p_scan.add_argument(
        '--fresh', action='store_true',
        help='Discard previous catalog and start fresh',
    )

    p_cls = sub.add_parser('classify', help='Batch classify with sonnet')
    p_cls.add_argument(
        '--model', default='sonnet', help='Claude model (default: sonnet)',
    )
    p_cls.add_argument(
        '--fresh', action='store_true', help='Re-classify all files',
    )

    p_rpt = sub.add_parser('report', help='Generate triage report')

    p_exec = sub.add_parser('execute', help='Move files per plan')
    p_exec.add_argument('path', help='OneDrive root path')
    p_exec.add_argument(
        '--dry-run', action='store_true', help='Preview only',
    )
    p_exec.add_argument(
        '--confirm', action='store_true', help='Actually move files',
    )

    args = parser.parse_args()
    cmds = {
        'scan': cmd_scan,
        'classify': cmd_classify,
        'report': cmd_report,
        'execute': cmd_execute,
    }
    fn = cmds.get(args.command)
    if fn:
        fn(args)
    else:
        parser.print_help()


if __name__ == '__main__':
    main()
